from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    prs = Presentation()

    # Title Slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "PhonePe Pulse Data Visualization & Analysis"
    subtitle.text = "LabMentix Data Science Internship Project\nDeveloped by Nikhil Kumar"

    # Slide 1: Introduction
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Introduction & Objective"
    
    content = slide.placeholders[1]
    content.text = "Goal: To provide clear insights into PhonePe transaction data (2018-2024)."

    # Slide 2: Technologies
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Technologies Used"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Core Stack:"
    p = tf.add_paragraph()
    p.text = "- Python, Pandas, SQLite"
    p = tf.add_paragraph()
    p.text = "- Streamlit for Frontend/Dashboard"
    p = tf.add_paragraph()
    p.text = "- Plotly for Interactive Graphics"
    p = tf.add_paragraph()
    p.text = "- Scikit-learn for ML Forecasting"

    # Slide 3: ETL Pipeline
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "ETL Pipeline Orchestration"
    content = slide.placeholders[1]
    content.text = "1. Extract: Git-cloned raw PhonePe repository.\n2. Transform: JSON processing to CSV/DF.\n3. Load: SQLite database integration (phonepe_pulse.db)."

    # Slide 4: Key Features
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Interactive Dashboard Features"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Dashboard Modules:"
    p = tf.add_paragraph()
    p.text = "- 🌍 Geographical Heatmaps"
    p = tf.add_paragraph()
    p.text = "- 📊 Yearly & Quarterly Trends"
    p = tf.add_paragraph()
    p.text = "- 🤝 Insurance Analysis"
    p = tf.add_paragraph()
    p.text = "- 🔮 ML-powered Growth Predictions"

    # Slide 5: ML Component
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Growth Forecasting (ML)"
    content = slide.placeholders[1]
    content.text = "Using a Random Forest Regressor to predict Transaction Amount for any State/Quarter/Year combo."

    # Slide 6: Deployment
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Global Deployment"
    content = slide.placeholders[1]
    content.text = "Hosted on: Streamlit Community Cloud"
    p = content.text_frame.add_paragraph()
    p.text = "Access URL: https://viraj357-phonepay-app-jokpqy.streamlit.app/"

    # Slide 7: Conclusion
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = "Conclusion"
    content = slide.placeholders[1]
    content.text = "This project demonstrates the end-to-end data lifecycle: from raw data extraction to interactive web visualization."

    # Save
    prs.save('PhonePe_Pulse_Presentation.pptx')
    print("Presentation generated: PhonePe_Pulse_Presentation.pptx")

if __name__ == "__main__":
    create_presentation()
